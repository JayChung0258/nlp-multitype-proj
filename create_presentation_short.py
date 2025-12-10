"""
Create 10-minute PowerPoint presentation for T1/T3 improvement experiments.
~10-12 slides for 10-minute presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

print("Creating 10-minute PowerPoint presentation...")
print("="*70)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add bullet point slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add slide with image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add image
    if Path(image_path).exists():
        left = Inches(0.8)
        top = Inches(1.3)
        slide.shapes.add_picture(str(image_path), left, top, width=Inches(8.4))
    
    # Add caption if provided
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.6))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_frame.paragraphs[0].font.size = Pt(16)
        caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        caption_frame.paragraphs[0].font.bold = True
    
    return slide

# ============================================================================
# Slide 1: Title (30 seconds)
# ============================================================================
print("[1/11] Creating title slide...")
add_title_slide(prs, 
                "Improving T1 & T3 Classification",
                "Advanced Techniques for Hard Classes\n\nJay, Yixuan, Sining\nDecember 2025")

# ============================================================================
# Slide 2: Problem Statement (45 seconds)
# ============================================================================
print("[2/11] Creating problem slide...")
add_content_slide(prs, "The Challenge",
                  [
                      "Baseline: DeBERTa-v3-base (F1: 0.711)",
                      "",
                      "  T1 (Human Original):      0.505  ← Low",
                      "  T2 (LLM Generated):       0.916  ✓",
                      "  T3 (Human Paraphrase):    0.515  ← Lowest",
                      "  T4 (LLM Paraphrase):      0.907  ✓",
                      "",
                      "Why Hard? T1 & T3 both human-written",
                      "Only differ by synonym substitution",
                      "",
                      "Goal: Improve T1 & T3 without breaking T2/T4"
                  ])

# ============================================================================
# Slide 3: Three Methods Tested (1 minute)
# ============================================================================
print("[3/11] Creating methods overview...")
add_content_slide(prs, "Three Improvement Strategies",
                  [
                      "1. Class Weighting",
                      "   Penalize T3 errors more heavily",
                      "",
                      "2. LoRA Fine-Tuning",
                      "   Parameter-efficient training (0.3% weights)",
                      "",
                      "3. Ensemble Voting",
                      "   Combine DeBERTa + RoBERTa + BERT",
                      "",
                      "",
                      "All methods designed to boost T1/T3 performance"
                  ])

# ============================================================================
# Slide 4: Method 1 - Class Weights (1.5 minutes)
# ============================================================================
print("[4/11] Creating class weights slide...")
add_content_slide(prs, "Method 1: Class Weighting",
                  [
                      "Approach: Increase T3 loss penalty by 1.5×",
                      "",
                      "Results:",
                      "  T3: 0.515 → 0.667  (+15%)  ✓",
                      "  T1: 0.505 → 0.060  (-88%)  ✗✗",
                      "  Overall: 0.711 → 0.640  (-10%)",
                      "",
                      "Why Failed:",
                      "  Model \"gamed\" the loss function",
                      "  Sacrificed T1 entirely to optimize T3",
                      "",
                      "Only 45/750 T1 samples classified correctly!"
                  ])

# ============================================================================
# Slide 5: Class Weights Visualization (1 minute)
# ============================================================================
print("[5/11] Creating class weights confusion matrix...")
add_image_slide(prs, "Class Weights: T1 Catastrophic Collapse",
                "visualizations/confusion_matrix_class_weights.png",
                "T1 completely abandoned - only 6% correct!")

# ============================================================================
# Slide 6: Method 2 - LoRA (1.5 minutes)
# ============================================================================
print("[6/11] Creating LoRA results slide...")
add_content_slide(prs, "Method 2: LoRA Fine-Tuning",
                  [
                      "Approach: Train only 0.3% of parameters (rank=8)",
                      "",
                      "Results: Catastrophic failure across ALL classes",
                      "  T1: 0.505 → 0.461",
                      "  T2: 0.916 → 0.471  (-45%!)  ✗✗",
                      "  T3: 0.515 → 0.524  (+1%)",
                      "  T4: 0.907 → 0.700  (-21%)",
                      "  Overall: 0.711 → 0.539  (-24%)",
                      "",
                      "Why Failed: Insufficient model capacity",
                      "  4-way classification too complex for 0.3% weights"
                  ])

# ============================================================================
# Slide 7: Method 3 - Ensemble Strategy (1 minute)
# ============================================================================
print("[7/11] Creating ensemble strategy slide...")
add_content_slide(prs, "Method 3: Ensemble Voting",
                  [
                      "Hypothesis: Diverse models correct each other",
                      "",
                      "Individual Model Weaknesses:",
                      "  DeBERTa:  Balanced (0.711 F1)",
                      "  RoBERTa:  T3 collapsed (0.125 F1)  ✗",
                      "  BERT:     T1 collapsed (0.307 F1)  ✗",
                      "",
                      "Strategy: Majority voting",
                      "  Each model predicts independently",
                      "  Take most common prediction",
                      "",
                      "Expected: Diversity should fix weaknesses"
                  ])

# ============================================================================
# Slide 8: Ensemble Results (1 minute)
# ============================================================================
print("[8/11] Creating ensemble results slide...")
add_content_slide(prs, "Ensemble Results: Also Failed",
                  [
                      "              DeBERTa   Ensemble   Change",
                      "Overall:       0.711     0.705    -0.006  ✗",
                      "T1:            0.505     0.551    +0.046  ✓",
                      "T3:            0.515     0.472    -0.043  ✗✗",
                      "",
                      "Ensemble WORSE than single DeBERTa!",
                      "",
                      "Why Failed:",
                      "  RoBERTa's extreme T3 failure (12.5%)",
                      "  polluted the voting",
                      "",
                      "409/747 T3 samples misclassified as T1"
                  ])

# ============================================================================
# Slide 9: Visual Comparison - All Methods (1.5 minutes)
# ============================================================================
print("[9/11] Creating comprehensive comparison...")
add_image_slide(prs, "All Methods Comparison",
                "visualizations/model_comparison.png",
                "DeBERTa baseline outperforms all advanced techniques")

# ============================================================================
# Slide 10: Key Learnings (1 minute)
# ============================================================================
print("[10/11] Creating key learnings slide...")
add_content_slide(prs, "Key Learnings",
                  [
                      "1. Class Weights Harmful on Balanced Data",
                      "   Created pathological optimization trade-offs",
                      "",
                      "2. LoRA Insufficient for Complex Tasks",
                      "   0.3% parameters too constrained for 4-way task",
                      "",
                      "3. Bad Models Poison Ensembles",
                      "   Quality > Diversity in ensemble design",
                      "",
                      "Universal Lesson:",
                      "  Simple methods often beat complex tricks",
                      "  when assumptions don't align with problem"
                  ])

# ============================================================================
# Slide 11: Conclusions (1 minute)
# ============================================================================
print("[11/11] Creating conclusions slide...")
add_content_slide(prs, "Conclusions",
                  [
                      "Best Model: DeBERTa-v3-base (no tricks)",
                      "  Macro-F1: 0.711",
                      "  T1: 0.505  |  T3: 0.515",
                      "",
                      "All advanced techniques failed:",
                      "  Class weights: T1 collapsed (-88%)",
                      "  LoRA: All classes degraded (-24%)",
                      "  Ensemble: Worse than single model",
                      "",
                      "Why T3 is Hard:",
                      "  Human paraphrases ≈ Human originals",
                      "  Only differ by synonyms",
                      "  Fundamentally challenging classification",
                      "",
                      "Current 0.515 F1 is strong performance!"
                  ])

# ============================================================================
# Save presentation
# ============================================================================
output_file = "T1_T3_Improvement_10min.pptx"
prs.save(output_file)

print("\n" + "="*70)
print(f"✓ 10-minute PowerPoint created successfully!")
print(f"  Saved as: {output_file}")
print("="*70)
print("\nPresentation structure:")
print("  • 11 slides total")
print("  • ~1 minute per slide")
print("  • 3 key visualizations embedded")
print("  • Clear narrative arc")
print("\nTiming breakdown:")
print("  Slides 1-2:   Problem (1:15)")
print("  Slide 3:      Methods overview (1:00)")
print("  Slides 4-5:   Class weights (2:30)")
print("  Slide 6:      LoRA (1:30)")
print("  Slides 7-8:   Ensemble (2:00)")
print("  Slide 9:      Comparison (1:30)")
print("  Slides 10-11: Learnings & conclusions (2:00)")
print("  ─────────────────────────────")
print("  Total: ~10 minutes")
print("="*70)
